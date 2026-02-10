
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Define corporate colors
    # TechSOS Blue: 0, 80, 156 (approx)
    TITLE_COLOR = RGBColor(0, 51, 102)
    TEXT_COLOR = RGBColor(60, 60, 60)

    def add_slide(layout_idx, title_text, content_text=None):
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
        title.text_frame.paragraphs[0].font.bold = True
        
        # Content
        if content_text:
            # Depending on layout, content placeholder might vary
            if len(slide.placeholders) > 1:
                content = slide.placeholders[1]
                content.text = content_text
        
        return slide

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0]) # Title Slide
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    # Add Logo
    logo_path = r"C:/Users/Marco/.gemini/antigravity/brain/cebdbbb7-ba52-4d53-9f51-90b8729b31eb/techsos_logo_v7_fixed_handle_1769175263786.png"
    try:
        slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.5), height=Inches(2.0))
    except Exception as e:
        print(f"Could not load logo: {e}")

    title.text = "TechSOS"
    subtitle.text = "Plano de Negócios & Investimento\nAlvalade, Lisboa - 2026"
    
    # Customizing Title
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    title.text_frame.paragraphs[0].font.size = Pt(54)

    # Slide 2: Oportunidade
    s2 = add_slide(1, "A Oportunidade", 
                   "Aquisição e modernização da antiga 'Conficell' em Alvalade.\n\n"
                   "• Localização Premium (Av. da Igreja)\n"
                   "• Faturação Validada de 23.000€/mês\n"
                   "• Investimento de Baixo Risco (ROI < 8 Meses)")

    # Slide 3: O Upgrade Digital
    s3 = add_slide(1, "Diferenciação: O Upgrade Digital",
                   "Transformamos uma loja tradicional num 'Tech Hub':\n\n"
                   "✅ Site E-commerce com Stock Real (Shopify)\n"
                   "✅ Assistente AI 24/7\n"
                   "✅ Pagamentos c/ Klarna (Pague em 3x)\n"
                   "✅ Reparações de Nível 3 (Placa)")

    # Slide 4: Micro-Localização
    s4 = add_slide(1, "Localização Imbatível",
                   "Avenida da Igreja, 15 - Alvalade\n\n"
                   "• Âncoras: McDonald's, Bancos, Pingo Doce\n"
                   "• Tráfego: Um dos bairros mais movimentados de Lisboa\n"
                   "• Público: Estudantes (Cidade Universitária) e Residentes Premium")

    # Slide 5: A Fórmula do Lucro (Margens)
    s5 = add_slide(1, "Margens Explosivas (Nível 3)",
                   "Ao internalizar reparações complexas, maximizamos o lucro:\n\n"
                   "• Custo Ecrã iPhone 11: 17€\n"
                   "• Preço Venda Mercado: 85€ - 100€\n"
                   "👉 Margem de 400% por serviço\n\n"
                   "Não vendemos apenas peças, vendemos especialização.")

    # Slide 6: O Espelho da Realidade (Validação)
    s6 = add_slide(1, "Prova de Conceito",
                   "Dados reais da gestão anterior (Out/Nov 2025):\n\n"
                   "💰 Faturação Mensal: ~23.700€\n"
                   "📈 Ticket Médio Elevado (> 1.300€)\n\n"
                   "O mercado existe e já compra nesta loja.")

    # Slide 7: Investimento Incial
    s7 = add_slide(1, "Investimento Necessário",
                   "Valor Total: ~14.300€\n\n"
                   "1. Trespasse (Direito Comercial): 5.000€\n"
                   "2. Stock Dispositivos (Arranque): 5.000€\n"
                   "3. Stock Existente: 1.500€\n"
                   "4. Cauções: 2.300€\n\n"
                   "Payback estimado em menos de 8 meses.")

    # Slide 8: Projeções (Ano 1)
    s8 = add_slide(1, "Projeção Financeira (Ano 1)",
                   "Cenário Conservador (18k/mês):\n\n"
                   "• Vendas Totais: 216.000€\n"
                   "• Lucro Bruto: 64.800€\n"
                   "• Custos Fixos (c/ Salários): 60.600€\n\n"
                   "👉 Lucro Líquido Positivo desde o Dia 1.\n"
                   "(Pagando Salário Técnico + Atendimento)")

    # Slide 9: Roadmap
    s9 = add_slide(1, "Roadmap de Crescimento",
                   "Mês 1: Ativação Klarna e Site\n"
                   "Mês 3: SEO Local 'Reparação iPhone Alvalade'\n"
                   "Mês 6: Contratação Técnico Júnior\n"
                   "Ano 2: Expansão Vendas Online Nacional")

    # Slide 10: Conclusão
    s10 = add_slide(0, "A Proposta Final")
    title = s10.shapes.title
    title.text = "Invista na TechSOS"
    subtitle = s10.placeholders[1]
    subtitle.text = "Negócio Validado + Gestão Moderna\n\nVamos fechar negócio?"

    output_path = r"c:\Users\Marco\.gemini\antigravity\brain\cebdbbb7-ba52-4d53-9f51-90b8729b31eb\Apresentacao_TechSOS.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
